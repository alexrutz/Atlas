"""
Datei-Parser für verschiedene Dokumentformate.

Unterstützt: PDF, DOCX, XLSX, PPTX, TXT, MD, CSV, HTML
"""

import logging
from dataclasses import dataclass, field
from pathlib import Path

logger = logging.getLogger(__name__)


@dataclass
class ParsedSection:
    """Ein Abschnitt aus einem geparsten Dokument."""
    header: str | None
    content: str
    page_number: int | None = None


@dataclass
class ParsedDocument:
    """Ergebnis des Dokument-Parsings."""
    text: str                                  # Gesamter extrahierter Text
    sections: list[ParsedSection] = field(default_factory=list)  # Erkannte Abschnitte
    page_count: int | None = None              # Seitenanzahl (bei PDFs)
    metadata: dict = field(default_factory=dict)


def parse_document(file_path: str, file_type: str) -> ParsedDocument:
    """
    Parst ein Dokument und extrahiert den Text.

    Args:
        file_path: Pfad zur Datei
        file_type: Dateityp (z.B. '.pdf', '.docx')

    Returns:
        ParsedDocument mit extrahiertem Text und Abschnitten
    """
    parsers = {
        ".pdf": _parse_pdf,
        ".docx": _parse_docx,
        ".doc": _parse_docx,
        ".xlsx": _parse_xlsx,
        ".xls": _parse_xlsx,
        ".pptx": _parse_pptx,
        ".txt": _parse_text,
        ".md": _parse_text,
        ".csv": _parse_csv,
        ".html": _parse_html,
        ".xml": _parse_html,
        ".json": _parse_text,
    }

    parser = parsers.get(file_type.lower())
    if not parser:
        raise ValueError(f"Nicht unterstütztes Dateiformat: {file_type}")

    return parser(file_path)


def _parse_pdf(file_path: str) -> ParsedDocument:
    """Parst ein PDF-Dokument.

    Wenn vlm_always aktiviert ist, wird VLM-OCR für ALLE PDFs verwendet
    (bessere Strukturerkennung bei Tabellen, Spalten, Überschriften).
    Andernfalls wird VLM/Tesseract nur als Fallback für gescannte PDFs genutzt.
    """
    from pypdf import PdfReader
    from app.core.config import settings

    reader = PdfReader(file_path)
    page_count = len(reader.pages)

    # --- VLM-always-Modus: VLM für alle PDFs (layout-aware Chunks) ---
    use_vlm_always = (
        settings.documents.vlm_always
        and settings.documents.ocr_backend == "vlm"
        and settings.vlm_ocr.enabled
    )
    if use_vlm_always:
        logger.info(f"VLM-always: Verarbeite PDF mit Layout-aware OCR: {file_path}")
        vlm_sections, vlm_text = _vlm_ocr_pdf(file_path)
        if vlm_text:
            return ParsedDocument(
                text="\n\n".join(vlm_text),
                sections=vlm_sections,
                page_count=page_count,
            )
        logger.warning(f"VLM-always lieferte keinen Text, Fallback auf pypdf: {file_path}")

    # --- Standard-Modus: Text mit pypdf extrahieren ---
    sections = []
    full_text = []

    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            sections.append(ParsedSection(
                header=None,
                content=text,
                page_number=i + 1,
            ))
            full_text.append(text)

    # Wenn kein Text extrahiert wurde (z.B. gescannte Dokumente), OCR verwenden
    if not full_text:
        if settings.documents.ocr_enabled:
            logger.info(f"Kein Text in PDF gefunden, starte OCR für: {file_path}")

            backend = getattr(settings.documents, "ocr_backend", "tesseract")
            if backend == "vlm" and settings.vlm_ocr.enabled:
                ocr_sections, ocr_text = _vlm_ocr_pdf(file_path)
            else:
                ocr_sections, ocr_text = _ocr_pdf(file_path, settings.documents.ocr_language)

            if ocr_text:
                return ParsedDocument(
                    text="\n\n".join(ocr_text),
                    sections=ocr_sections,
                    page_count=page_count,
                )
            logger.warning(f"OCR konnte keinen Text extrahieren: {file_path}")
        else:
            logger.warning(f"Kein Text in PDF und OCR deaktiviert: {file_path}")

    return ParsedDocument(
        text="\n\n".join(full_text),
        sections=sections,
        page_count=page_count,
    )


def _clean_ocr_text(text: str) -> str:
    """Bereinigt OCR-Text von typischen Artefakten.

    Entfernt:
    - Zeilen die überwiegend aus Sonderzeichen bestehen
    - Übermäßige Leerzeichen und Steuerzeichen
    - Sehr kurze Zeilen die wahrscheinlich Rauschen sind
    """
    import re

    lines = text.split("\n")
    cleaned = []
    for line in lines:
        line = line.strip()
        if not line:
            continue
        # Anteil alphanumerischer Zeichen berechnen
        alnum_count = sum(1 for c in line if c.isalnum())
        if len(line) > 0 and alnum_count / len(line) < 0.3:
            # Zeile besteht zu >70% aus Sonderzeichen → Rauschen
            continue
        if len(line) < 3:
            continue
        # Mehrfache Leerzeichen zusammenfassen
        line = re.sub(r" {2,}", " ", line)
        cleaned.append(line)

    return "\n".join(cleaned)


def _ocr_pdf(file_path: str, ocr_language: str = "deu+eng") -> tuple[list[ParsedSection], list[str]]:
    """Führt OCR auf allen Seiten eines PDFs durch."""
    try:
        from pdf2image import convert_from_path
        import pytesseract
    except ImportError:
        logger.error("pdf2image oder pytesseract nicht installiert, OCR nicht möglich")
        return [], []

    sections = []
    full_text = []

    try:
        images = convert_from_path(file_path, dpi=300)
        for i, image in enumerate(images):
            raw_text = pytesseract.image_to_string(image, lang=ocr_language)
            text = _clean_ocr_text(raw_text)
            if text.strip() and len(text.strip()) >= 20:
                sections.append(ParsedSection(
                    header=None,
                    content=text.strip(),
                    page_number=i + 1,
                ))
                full_text.append(text.strip())
            elif text.strip():
                logger.debug(f"OCR-Seite {i+1} übersprungen (zu wenig brauchbarer Text: {len(text.strip())} Zeichen)")
    except Exception as e:
        logger.error(f"OCR-Fehler: {e}")

    return sections, full_text


def _vlm_ocr_pdf(file_path: str) -> tuple[list[ParsedSection], list[str]]:
    """Führt VLM-OCR (Layout-as-thought) auf allen Seiten eines PDFs durch.

    Verwendet den VlmOcrService asynchron. Da file_parsers synchron aufgerufen
    wird, wird ein neuer Event-Loop gestartet falls nötig.
    """
    import asyncio

    try:
        from app.services.vlm_ocr_service import VlmOcrService
        service = VlmOcrService()

        try:
            loop = asyncio.get_running_loop()
        except RuntimeError:
            loop = None

        if loop and loop.is_running():
            # Bereits in einem async-Kontext → neuen Thread nutzen
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
                future = pool.submit(asyncio.run, service.ocr_pdf_pages(file_path))
                return future.result()
        else:
            return asyncio.run(service.ocr_pdf_pages(file_path))

    except Exception as e:
        logger.error(f"VLM-OCR fehlgeschlagen, Fallback auf Tesseract: {e}")
        from app.core.config import settings
        return _ocr_pdf(file_path, settings.documents.ocr_language)


def _parse_docx(file_path: str) -> ParsedDocument:
    """Parst ein Word-Dokument."""
    from docx import Document as DocxDocument

    doc = DocxDocument(file_path)
    sections = []
    current_section = None
    current_text = []

    for para in doc.paragraphs:
        if para.style.name.startswith("Heading"):
            if current_text:
                sections.append(ParsedSection(
                    header=current_section, content="\n".join(current_text),
                ))
            current_section = para.text
            current_text = []
        else:
            if para.text.strip():
                current_text.append(para.text)

    if current_text:
        sections.append(ParsedSection(header=current_section, content="\n".join(current_text)))

    full_text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return ParsedDocument(text=full_text, sections=sections)


def _parse_xlsx(file_path: str) -> ParsedDocument:
    """Parst eine Excel-Datei."""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, data_only=True)
    sections = []
    full_text = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        rows = []
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                rows.append(row_text)

        if rows:
            section_text = "\n".join(rows)
            sections.append(ParsedSection(header=f"Blatt: {sheet_name}", content=section_text))
            full_text.append(f"[{sheet_name}]\n{section_text}")

    return ParsedDocument(text="\n\n".join(full_text), sections=sections)


def _parse_pptx(file_path: str) -> ParsedDocument:
    """Parst eine PowerPoint-Datei."""
    from pptx import Presentation

    prs = Presentation(file_path)
    sections = []
    full_text = []

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)

        if texts:
            slide_text = "\n".join(texts)
            sections.append(ParsedSection(header=f"Folie {i + 1}", content=slide_text, page_number=i + 1))
            full_text.append(slide_text)

    return ParsedDocument(text="\n\n".join(full_text), sections=sections, page_count=len(prs.slides))


def _parse_text(file_path: str) -> ParsedDocument:
    """Parst eine Textdatei."""
    with open(file_path, encoding="utf-8", errors="replace") as f:
        text = f.read()
    return ParsedDocument(text=text, sections=[ParsedSection(header=None, content=text)])


def _parse_csv(file_path: str) -> ParsedDocument:
    """Parst eine CSV-Datei."""
    import csv
    rows = []
    with open(file_path, encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append(" | ".join(row))

    text = "\n".join(rows)
    return ParsedDocument(text=text, sections=[ParsedSection(header=None, content=text)])


def _parse_html(file_path: str) -> ParsedDocument:
    """Parst eine HTML/XML-Datei."""
    from bs4 import BeautifulSoup

    with open(file_path, encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "html.parser")

    text = soup.get_text(separator="\n", strip=True)
    return ParsedDocument(text=text, sections=[ParsedSection(header=None, content=text)])
